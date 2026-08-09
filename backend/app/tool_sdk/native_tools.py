from __future__ import annotations
import base64, hashlib, ipaddress, json, os, re, socket
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path
from urllib.parse import urljoin, urlparse
import bleach  # type: ignore[import-untyped]
import httpx
import sqlglot
from docx import Document
from openpyxl import load_workbook  # type: ignore[import-untyped]
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy import create_engine, text
from app.contracts.tool import Tool
from app.contracts.tool_models import ToolMetadata, ToolResult
from app.database.models.native_tool import (
    NativeConnection,
    NativeFile,
    NativeFileContent,
    NativeNotification,
)
from app.tool_sdk.errors import ToolSDKError, UnsafeOperationError


class NativeError(ToolSDKError):
    pass


class UnsafeQuery(NativeError):
    code, status_code = "UNSAFE_QUERY_REJECTED", 403


class UnsafeURL(NativeError):
    code, status_code = "UNSAFE_URL_REJECTED", 403


class ApprovalRequired(NativeError):
    code, status_code = "APPROVAL_REQUIRED", 409


ALLOWED_EXT = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".json": "application/json",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def metadata(
    name,
    display,
    description,
    permission,
    props,
    required=(),
    risk="read",
    timeout=30,
    idempotent=True,
):
    return ToolMetadata(
        name=name,
        display_name=display,
        description=description,
        category="native_enterprise",
        provider="native",
        version="1.0.0",
        permissions=(permission,),
        tags=("native", name.split("_")[0]),
        parameters={
            "type": "object",
            "properties": props,
            "required": list(required),
            "additionalProperties": False,
        },
        risk_level=risk,
        timeout_seconds=timeout,
        idempotent=idempotent,
        configuration_requirements=(name.split("_")[0],),
    )


def storage_root():
    root = Path(
        os.getenv(
            "NATIVE_FILE_STORAGE_ROOT",
            Path(__file__).parents[3] / "data" / "native-files",
        )
    ).resolve()
    root.mkdir(parents=True, exist_ok=True)
    return root


def detect(data, ext):
    if data.startswith(b"%PDF-"):
        return "application/pdf"
    if data.startswith(b"PK\x03\x04") and ext in {".docx", ".xlsx", ".pptx"}:
        return ALLOWED_EXT[ext]
    if b"\x00" in data[:4096]:
        return "application/octet-stream"
    return ALLOWED_EXT.get(ext)


def extract_bytes(data, ext):
    if ext in {".txt", ".md", ".csv", ".json"}:
        return data.decode("utf-8"), None
    if ext == ".pdf":
        reader = PdfReader(BytesIO(data))
        return "\n\n".join((p.extract_text() or "") for p in reader.pages), len(
            reader.pages
        )
    if ext == ".docx":
        return "\n".join(p.text for p in Document(BytesIO(data)).paragraphs), None
    if ext == ".xlsx":
        wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
        return "\n".join(
            ",".join("" if v is None else str(v) for v in row)
            for ws in wb.worksheets
            for row in ws.iter_rows(values_only=True)
        ), None
    if ext == ".pptx":
        deck = Presentation(BytesIO(data))
        return "\n".join(
            shape.text
            for slide in deck.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        ), len(deck.slides)
    raise NativeError("Unsupported file type")


class FileUploadTool(Tool):
    metadata = metadata(
        "file_upload",
        "File Upload",
        "Validate, scan, store, extract, and index an enterprise document",
        "files.upload",
        {
            "filename": {"type": "string", "minLength": 1, "maxLength": 255},
            "content_base64": {
                "type": "string",
                "minLength": 1,
                "maxLength": 8_000_000,
            },
        },
        ("filename", "content_base64"),
        "write",
        60,
    )

    async def execute(self, input_data, context):
        name = Path(input_data["filename"]).name
        ext = Path(name).suffix.lower()
        if ext not in ALLOWED_EXT or name != input_data["filename"]:
            raise NativeError("Unsupported or unsafe filename")
        try:
            data = base64.b64decode(input_data["content_base64"], validate=True)
        except Exception as exc:
            raise NativeError("Upload content is not valid base64") from exc
        limit = int(os.getenv("NATIVE_FILE_MAX_BYTES", "5242880"))
        if len(data) > limit:
            raise NativeError("File exceeds configured size limit")
        mime = detect(data, ext)
        if mime != ALLOWED_EXT[ext]:
            raise NativeError("File content does not match its extension")
        if b"EICAR-STANDARD-ANTIVIRUS-TEST-FILE" in data:
            raise NativeError("File was quarantined by the malware scanner")
        checksum = hashlib.sha256(data).hexdigest()
        db = context.db_session
        existing = (
            db.query(NativeFile)
            .filter_by(tenant_id=context.tenant_id, checksum=checksum)
            .first()
        )
        if existing:
            return ToolResult.succeeded(file_item(existing) | {"duplicate": True})
        safe = re.sub(r"[^A-Za-z0-9._-]", "_", name)
        key = f"{context.tenant_id}/{checksum[:2]}/{checksum}-{safe}"
        path = (storage_root() / key).resolve()
        if storage_root() not in path.parents:
            raise UnsafeOperationError("Unsafe storage path")
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(data)
        row = NativeFile(
            tenant_id=context.tenant_id,
            original_filename=name,
            normalized_filename=safe,
            object_key=key,
            mime_type=mime,
            extension=ext,
            byte_size=len(data),
            checksum=checksum,
            owner_id=context.actor_id,
            scan_status="safe",
            processing_status="processing",
        )
        db.add(row)
        db.commit()
        db.refresh(row)
        try:
            content, pages = extract_bytes(data, ext)
            content = content[:2_000_000]
            db.add(
                NativeFileContent(
                    file_id=row.id,
                    tenant_id=context.tenant_id,
                    sequence=0,
                    text=content,
                    character_count=len(content),
                    metadata_json={},
                )
            )
            row.processing_status = "indexed"
            row.page_count = pages
            row.extractor_version = "1.0.0"
            db.commit()
        except Exception:
            row.processing_status = "failed"
            row.error_code = "FILE_PROCESSING_FAILED"
            db.commit()
            raise
        return ToolResult.succeeded(file_item(row))


def file_item(x):
    return {
        "id": x.id,
        "filename": x.original_filename,
        "mime_type": x.mime_type,
        "size": x.byte_size,
        "checksum": x.checksum,
        "scan_status": x.scan_status,
        "status": x.processing_status,
        "page_count": x.page_count,
        "created_at": x.created_at.isoformat(),
    }


class FileOperationTool(Tool):
    def __init__(self, op):
        self.op = op
        permission = {
            "read": "files.read",
            "extract": "files.extract",
            "summarize": "files.summarize",
            "search": "files.search",
        }[op]
        props = {
            "file_id": {"type": "string", "maxLength": 36},
            "query": {"type": "string", "maxLength": 500},
            "focus": {"type": "string", "maxLength": 500},
        }
        required = () if op == "search" else ("file_id",)
        self.metadata = metadata(
            f"file_{op}",
            f"File {op.title()}",
            f"Secure native document {op}",
            permission,
            props,
            required,
        )

    async def execute(self, input_data, context):
        db = context.db_session
        if self.op == "search":
            q = input_data.get("query", "").lower()
            rows = (
                db.query(NativeFileContent, NativeFile)
                .join(NativeFile, NativeFile.id == NativeFileContent.file_id)
                .filter(
                    NativeFileContent.tenant_id == context.tenant_id,
                    NativeFile.processing_status == "indexed",
                )
                .limit(50)
                .all()
            )
            items = [
                file_item(f)
                | {
                    "excerpt": c.text[
                        max(0, c.text.lower().find(q) - 80) : c.text.lower().find(q)
                        + 240
                    ]
                    if q in c.text.lower()
                    else c.text[:240]
                }
                for c, f in rows
                if not q or q in c.text.lower() or q in f.original_filename.lower()
            ]
            return ToolResult.succeeded({"items": items[:25], "count": len(items[:25])})
        row = (
            db.query(NativeFile)
            .filter_by(id=input_data["file_id"], tenant_id=context.tenant_id)
            .first()
        )
        if not row:
            raise NativeError("File not found")
        if row.scan_status != "safe":
            raise NativeError("File is quarantined")
        content = (
            db.query(NativeFileContent)
            .filter_by(file_id=row.id, tenant_id=context.tenant_id)
            .order_by(NativeFileContent.sequence)
            .all()
        )
        full = "\n".join(x.text for x in content)
        if self.op in {"read", "extract"}:
            return ToolResult.succeeded({"file": file_item(row), "content": full})
        sentences = re.split(r"(?<=[.!?])\s+", full)
        summary = " ".join(sentences[:8])[:4000]
        return ToolResult.succeeded(
            {
                "file_id": row.id,
                "summary": summary,
                "mode": "deterministic_extractive",
                "citations": [{"section": 0}],
            }
        )


FORBIDDEN_SQL = {
    "Insert",
    "Update",
    "Delete",
    "Drop",
    "Alter",
    "Create",
    "Truncate",
    "Command",
    "Merge",
    "Grant",
    "Revoke",
    "Copy",
}


def safe_sql(query, allowed_tables, max_rows):
    if ";" in query.rstrip(";") or "--" in query or "/*" in query:
        raise UnsafeQuery("Multiple statements and comments are not allowed")
    try:
        tree = sqlglot.parse_one(query, read="postgres")
    except Exception as exc:
        raise UnsafeQuery("Query could not be parsed") from exc
    if tree.key.upper() not in {"SELECT", "UNION"}:
        raise UnsafeQuery("Only SELECT queries are allowed")
    if any(node.__class__.__name__ in FORBIDDEN_SQL for node in tree.walk()):
        raise UnsafeQuery("Mutating SQL is not allowed")
    tables = {t.name for t in tree.find_all(sqlglot.exp.Table)}
    if not tables <= set(allowed_tables):
        raise UnsafeQuery("Query references a table outside the allowlist")
    if not tree.args.get("limit"):
        tree = tree.limit(max_rows)
    return tree.sql(dialect="postgres"), hashlib.sha256(
        tree.sql().encode()
    ).hexdigest()[:16]


class DatabaseQueryTool(Tool):
    metadata = metadata(
        "database_query",
        "Database Tool",
        "Execute validated read-only queries against approved connections",
        "database.query",
        {
            "connection_id": {"type": "string", "maxLength": 36},
            "query": {"type": "string", "minLength": 1, "maxLength": 10000},
            "row_limit": {
                "type": "integer",
                "minimum": 1,
                "maximum": 1000,
                "default": 100,
            },
        },
        ("connection_id", "query"),
        timeout=30,
    )

    async def execute(self, input_data, context):
        db = context.db_session
        conn = (
            db.query(NativeConnection)
            .filter_by(
                id=input_data["connection_id"],
                tenant_id=context.tenant_id,
                kind="database",
                enabled=True,
            )
            .first()
        )
        if not conn:
            raise NativeError("Database connection is not configured")
        config = conn.safe_config or {}
        query, fingerprint = safe_sql(
            input_data["query"],
            config.get("allowed_tables", []),
            min(input_data["row_limit"], config.get("max_rows", 100)),
        )
        url = os.getenv(f"NATIVE_DB_URL_{conn.id.replace('-', '_').upper()}")
        if not url:
            raise NativeError("Database secret reference is not resolved")
        engine = create_engine(url)
        start = datetime.now(UTC)
        with engine.connect() as c:
            if engine.dialect.name == "postgresql":
                c.execute(text("SET TRANSACTION READ ONLY"))
            result = c.execute(text(query))
            rows = [dict(x) for x in result.mappings().all()]
        return ToolResult.succeeded(
            {
                "columns": list(rows[0]) if rows else [],
                "rows": rows,
                "row_count": len(rows),
                "query_fingerprint": fingerprint,
                "duration_ms": (datetime.now(UTC) - start).total_seconds() * 1000,
                "connection": conn.display_name,
            }
        )


FORBIDDEN_HEADERS = {
    "authorization",
    "proxy-authorization",
    "host",
    "cookie",
    "set-cookie",
    "forwarded",
    "x-forwarded-for",
    "connection",
}


def validate_url(base, path, allowed_paths):
    parsed = urlparse(base)
    if parsed.scheme != "https" or parsed.username or parsed.password:
        raise UnsafeURL("Only credential-free HTTPS base URLs are allowed")
    try:
        ips = {
            ipaddress.ip_address(x[4][0])
            for x in socket.getaddrinfo(parsed.hostname, parsed.port or 443)
        }
    except OSError as exc:
        raise UnsafeURL("Connection host could not be resolved") from exc
    if any(
        ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_reserved
        for ip in ips
    ):
        raise UnsafeURL("Connection resolves to a forbidden network")
    if not any(re.fullmatch(pattern, path) for pattern in allowed_paths):
        raise UnsafeURL("Endpoint path is not approved")
    return urljoin(base.rstrip("/") + "/", path.lstrip("/"))


class RestRequestTool(Tool):
    metadata = metadata(
        "rest_api_request",
        "REST API Tool",
        "Call an approved enterprise REST connection",
        "api.invoke",
        {
            "connection_id": {"type": "string", "maxLength": 36},
            "method": {
                "type": "string",
                "enum": ["GET", "POST", "PUT", "PATCH", "DELETE"],
            },
            "path": {"type": "string", "maxLength": 1000},
            "query": {"type": "object"},
            "headers": {"type": "object"},
            "body": {"type": "object"},
        },
        ("connection_id", "method", "path"),
        "write",
        30,
        False,
    )

    async def execute(self, input_data, context):
        db = context.db_session
        conn = (
            db.query(NativeConnection)
            .filter_by(
                id=input_data["connection_id"],
                tenant_id=context.tenant_id,
                kind="rest",
                enabled=True,
            )
            .first()
        )
        if not conn:
            raise NativeError("REST connection is not configured")
        cfg = conn.safe_config or {}
        method = input_data["method"]
        if method not in cfg.get("allowed_methods", ["GET"]):
            raise NativeError("HTTP method is not allowed")
        if (
            method != "GET"
            and "api.invoke.write" not in context.permissions
            and "tools.admin" not in context.permissions
        ):
            raise ApprovalRequired(
                "Write REST requests require elevated permission and approval"
            )
        url = validate_url(
            conn.base_url, input_data["path"], cfg.get("allowed_paths", [r"/.*"])
        )
        headers = {
            k: v
            for k, v in input_data.get("headers", {}).items()
            if k.lower() not in FORBIDDEN_HEADERS
        }
        token = os.getenv(f"NATIVE_REST_TOKEN_{conn.id.replace('-', '_').upper()}")
        headers |= {"Authorization": f"Bearer {token}"} if token else {}
        async with httpx.AsyncClient(
            timeout=min(cfg.get("timeout", 20), 30), follow_redirects=False
        ) as client:
            r = await client.request(
                method,
                url,
                params=input_data.get("query"),
                headers=headers,
                json=input_data.get("body"),
            )
        safe_headers = {
            k: v
            for k, v in r.headers.items()
            if k.lower() not in FORBIDDEN_HEADERS
            and k.lower() in {"content-type", "etag", "x-request-id", "x-ms-request-id"}
        }
        raw = r.content[: cfg.get("max_response_bytes", 1_000_000)]
        try:
            data = json.loads(raw)
        except Exception:
            data = raw.decode("utf-8", "replace")
        return ToolResult.succeeded(
            {
                "status_code": r.status_code,
                "headers": safe_headers,
                "data": data,
                "truncated": len(r.content) > len(raw),
            },
            provider_request_id=safe_headers.get("x-request-id"),
        )


class NotificationTool(Tool):
    def __init__(self, channel):
        self.channel = channel
        permission = (
            f"notifications.{channel}.{'create' if channel == 'alert' else 'send'}"
        )
        self.metadata = metadata(
            f"notification_{channel}_{'create' if channel == 'alert' else 'send'}",
            f"Notification {channel.title()}",
            f"Create a governed {channel} notification",
            permission,
            {
                "recipients": {
                    "type": "array",
                    "items": {"type": "string"},
                    "maxItems": 20,
                },
                "subject": {"type": "string", "maxLength": 200},
                "message": {"type": "string", "minLength": 1, "maxLength": 10000},
                "severity": {"type": "string", "enum": ["info", "warning", "critical"]},
                "idempotency_key": {"type": "string", "minLength": 8, "maxLength": 160},
            },
            ("recipients", "message", "idempotency_key"),
            "write",
            30,
            False,
        )

    async def execute(self, input_data, context):
        db = context.db_session
        prior = (
            db.query(NativeNotification)
            .filter_by(
                tenant_id=context.tenant_id,
                channel=self.channel,
                idempotency_key=input_data["idempotency_key"],
            )
            .first()
        )
        if prior:
            return ToolResult.succeeded(notification_item(prior) | {"duplicate": True})
        recipients = input_data["recipients"]
        domains = {
            x.lower()
            for x in os.getenv("NOTIFICATION_ALLOWED_EMAIL_DOMAINS", "").split(",")
            if x
        }
        external = self.channel == "email" and any(
            "@" not in x or x.rsplit("@", 1)[1].lower() not in domains
            for x in recipients
        )
        approval = (
            self.channel == "teams"
            or external
            or len(recipients) > 5
            or input_data.get("severity") == "critical"
        )
        message = bleach.clean(input_data["message"], tags=[], strip=True)
        row = NativeNotification(
            tenant_id=context.tenant_id,
            channel=self.channel,
            actor_id=context.actor_id,
            recipient_summary={"count": len(recipients), "recipients": recipients},
            subject=input_data.get("subject"),
            safe_message=message,
            status="pending_approval" if approval else "sent",
            approval_state="pending" if approval else "not_required",
            provider_message_id=None if approval else f"native-{os.urandom(6).hex()}",
            idempotency_key=input_data["idempotency_key"],
            sent_at=None if approval else datetime.now(UTC),
        )
        db.add(row)
        db.commit()
        db.refresh(row)
        return ToolResult.succeeded(
            notification_item(row)
            | ({"error": {"code": "APPROVAL_REQUIRED"}} if approval else {})
        )


def notification_item(x):
    return {
        "id": x.id,
        "channel": x.channel,
        "status": x.status,
        "approval_state": x.approval_state,
        "recipient_summary": x.recipient_summary,
        "provider_message_id": x.provider_message_id,
        "created_at": x.created_at.isoformat(),
    }


def native_tools():
    return [
        FileUploadTool(),
        FileOperationTool("read"),
        FileOperationTool("extract"),
        FileOperationTool("summarize"),
        FileOperationTool("search"),
        DatabaseQueryTool(),
        RestRequestTool(),
        NotificationTool("email"),
        NotificationTool("teams"),
        NotificationTool("alert"),
    ]
